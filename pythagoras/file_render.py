
from nhansohoc.settings import TEMPLATES
import os
from django.conf import settings
from io import StringIO, BytesIO
from django.http import HttpResponse
from django.template.loader import get_template, render_to_string
import xhtml2pdf.pisa as pisa
from django.contrib.staticfiles import finders

import pptx
from bs4 import BeautifulSoup



from pprint import pprint
def link_callback(uri, rel):
        """
        Convert HTML URIs to absolute system paths so xhtml2pdf can access those
        resources
        """
        result = finders.find(uri)
        if result:
                if not isinstance(result, (list, tuple)):
                        result = [result]
                result = list(os.path.realpath(path) for path in result)
                path=result[0]
        else:
                sUrl = settings.STATIC_URL        # Typically /static/
                sRoot = settings.STATIC_ROOT      # Typically /home/userX/project_static/
                mUrl = settings.MEDIA_URL         # Typically /media/
                mRoot = settings.MEDIA_ROOT       # Typically /home/userX/project_static/media/

                if uri.startswith(mUrl):
                        path = os.path.join(mRoot, uri.replace(mUrl, ""))
                elif uri.startswith(sUrl):
                        path = os.path.join(sRoot, uri.replace(sUrl, ""))
                else:
                        return uri

        # make sure that file exists
        if not os.path.isfile(path):
                raise Exception(
                        'media URI must start with %s or %s' % (sUrl, mUrl)
                )
        return path


def render(path: str, params: dict):
    template = get_template(path)
    html = template.render(params)
    result =  BytesIO()
    options = {'encoding': "UTF-8"}

    pdf = pisa.pisaDocument(BytesIO(html.encode("UTF-8")), result)
    if not pdf.err:
        return HttpResponse(result.getvalue(), content_type='application/pdf')
    else:
        return HttpResponse("Error Rendering PDF", status=400)


def render_pdf_view(path: str, params: dict):

    template = get_template(path)
    html = template.render(params)
    # Create a Django response object, and specify content_type as pdf
    response = HttpResponse(content_type='application/pdf')
    # download
    response['Content-Disposition'] = "'attachment; filename='report.pdf'"
    # display
#     response['Content-Disposition'] = "filename='report.pdf'"
    # find the template and render it.

    template = get_template(path)

    html = template.render(params) 
    
    # create a pdf
    pisa_status = pisa.CreatePDF(html, dest=response, link_callback=link_callback, debug=1)
    
#     return HttpResponse('HERE')   
    # if error then show some funy view
    if pisa_status.err:
        return HttpResponse('We had some errors <pre>' + html + '</pre>')
    
    return HttpResponse(response.getvalue(), content_type='application/pdf')
#     return response


def render_wkhtml(path: str, params: dict):
    from wkhtmltopdf.views import PDFTemplateResponse

    template = get_template(path)
    html = template.render(params)

    response = PDFTemplateResponse(request=params['request'],
                                   template=path,
                                   filename="hello.pdf",
                                   context= params,
                                   show_content_in_browser=True,
                                   cmd_options={'margin-top': 50,},
                                   )
    return response

def render_pdfkit(path: str, params: dict):
    import pdfkit

    
    template = get_template(path)

    html = template.render(params) 
    options = {
        'page-size': 'Letter',
        'encoding': "UTF-8",
    }
    pdf = pdfkit.from_url("http://localhost:8000/", options)
    response = HttpResponse(pdf, content_type='application/pdf')
    response['Content-Disposition'] = "'attachment; filename='report.pdf'"
    return response      

def replace_text(replacements, shapes):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:        
        if shape.has_text_frame and len(shape.text) > 0:
            txt = shape.text
            for match, replacement in replacements.items():
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if(not(not paragraph.runs)):
                            paragraph.runs[0].text = whole_text
            txt_dest = shape.text

def render_pptx(params:dict):
    # load up pptx template
    file_path = settings.TEMPLATES[0]["DIRS"][0]
    prs = pptx.Presentation(file_path / "template_2.pptx")
    slide_list = prs.slides
    # get details
    details = params.get("details", {})
    # life path
    lifepath = details['lifepath']
    lifepath_meaning = details['lifepath']['meaning']
    r_lifepath_meaning = BeautifulSoup(lifepath_meaning, "html5lib").get_text()
    destiny_meaning = details['destinypath']['meaning']
    r_destiny_meaning = BeautifulSoup(destiny_meaning, "html5lib").get_text()
    power_meaning = details['destinypath']['meaning']
    r_power_meaning = BeautifulSoup(power_meaning, "html5lib").get_text()
    hearth_desire_meaning = details['hearthdesire']['meaning']
    r_hearth_desire_meaning = BeautifulSoup(hearth_desire_meaning, "html5lib").get_text()
    personality_meaning = details['personality']['meaning']
    r_personality_meaning = BeautifulSoup(personality_meaning, "html5lib").get_text()
    attitude_meaning = details['attitudepath']['meaning']
    r_attitude_meaning = BeautifulSoup(attitude_meaning, "html5lib").get_text()
    daypath_meaning = details['daypath']['meaning']
    r_daypath_meaning = BeautifulSoup(daypath_meaning, "html5lib").get_text()
    yearpath_meaning = details['yearpath']['meaning']
    r_yearpath_meaning = BeautifulSoup(yearpath_meaning, "html5lib").get_text()

    last_name = details['info']['last_name'] 
    first_name = details['info']['first_name']
    replaces = {
        '{{fullname}}': last_name + ' ' + first_name,
        '{{birthday}}': str(details['dob']),
        '{{life_path_meaning}}': r_lifepath_meaning,
        '{{lp_num}}': str(details['info']['life_path_number']),
        '{{lp_strong}}': BeautifulSoup(lifepath['lp_strong'], "html5lib").get_text(),
        '{{lp_challenge}}': BeautifulSoup(lifepath['lp_challenge'], "html5lib").get_text(),
        '{{lp_improve}}': BeautifulSoup(lifepath['lp_improve'], "html5lib").get_text(),
        '{{lp_environment}}': BeautifulSoup(lifepath['lp_environment'], "html5lib").get_text(),
        '{{dp_num}}': str(details['info']['destiny_number']),
        '{{destiny_meaning}}': r_destiny_meaning,
        '{{pw_num}}': str(details['info']['power_number']),
        '{{power_meaning}}': r_power_meaning,
        '{{hd_num}}': str(details['info']['hearth_desire_number']),
        '{{hearth_desire_meaning}}': r_hearth_desire_meaning,
        '{{pe_num}}': str(details['info']['personality_number']),
        '{{personality_meaning}}': r_personality_meaning,
        '{{att_num}}': str(details['info']['attitude_number']),
        '{{attitude_meaning}}': r_attitude_meaning,
        '{{bd_num}}': str(details['info']['birthdate_day']),
        '{{birthdate_day_meaning}}': r_daypath_meaning,
        '{{bd_yr_num}}': str(details['info']['birthdate_day']),
        '{{birthdate_year_meaning}}': r_yearpath_meaning,
        '{{miss_num}}': str(details['info']['full_name_missing_numbers']),
        '{{miss_num_1}}': '',
        '{{miss_num_2}}': '',
        '{{miss_num_3}}': '',
        '{{miss_num_4}}': '',
        '{{miss_num_1_meaning}}': '',
        '{{miss_num_2_meaning}}': '',
        '{{miss_num_3_meaning}}': '',
        '{{miss_num_4_meaning}}': '',
    }

    # missing path 

    missing_path = details['missingpath']
    for idx, miss in enumerate(missing_path):
        miss_num_key = '{{miss_num_' + str(idx+1) + '}}'        
        miss_num_meaning_key = '{{miss_num_' + str(idx+1) + '_meaning}}'
        miss_num_meaning =  BeautifulSoup(miss['meaning'], "html5lib").get_text()
        replaces[miss_num_key] = miss['model_number']
        replaces[miss_num_meaning_key] = miss_num_meaning

    # pyramyd
    pyramid_path = details['pyramidpath']
    pyramid_ages = details['info']['pyramid_ages']
    pyramid_numbers = details['info']['pyramid_numbers']
    for idx, pyr in enumerate(pyramid_path):
        pyr_num_key = '{{py_num_' + str(idx+1) + '}}'
        pyr_age_key = '{{py_age_' + str(idx+1) + '}}'
        pyr_meaning_key = '{{py_' + str(idx+1) + '_meaning}}'
        pyr_meaning = BeautifulSoup(pyr['meaning'], "html5lib").get_text()

        replaces[pyr_num_key] = str(pyramid_numbers[idx])
        replaces[pyr_age_key] = str(pyramid_ages[idx])
        replaces[pyr_meaning_key] = pyr_meaning

    # challenges
    challenge_path = details['challengepath']
    challenge_numbers = details['info']['challenge_numbers']
    for idx, chall in enumerate(challenge_path):
        cll_num_key = '{{cll_num_' + str(idx+1) + '}}'        
        cll_meaning_key = '{{challenge_' + str(idx+1) + '_meaning}}'
        cll_meaning = BeautifulSoup(chall['meaning'], "html5lib").get_text()

        replaces[cll_num_key] = str(challenge_numbers[idx])
        replaces[cll_meaning_key] = cll_meaning


    shapes = []
    for slide in slide_list:
        for shape in slide.shapes:
            shapes.append(shape)                
    replace_text(replaces, shapes)
                    

    response = HttpResponse(content_type='application/vnd.ms-powerpoint')
    response['Content-Disposition'] = f'attachment; filename=report_{last_name}_{first_name}.pptx'
    source_stream = BytesIO()
    prs.save(source_stream)
    ppt = source_stream.getvalue()
    source_stream.close()
    response.write(ppt)

    return response
